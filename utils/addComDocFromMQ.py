#!/usr/bin/env python3
"""
addComDocFromMQ.py — RabbitMQ 消费者：处理 addCompanyDoc 任务

流程：
  1. 从 RabbitMQ 消费消息（addCompanyDoc_queue）
  2. 在 maindir/<companyId>/ 下创建企业 wiki 目录结构
  3. 拉取文档内容（URL 或本地路径）
  4. 调用 LLM 生成 wiki 页面并写入文件

用法：
    python utils/addComDocFromMQ.py            # 持续消费（Ctrl+C 退出）
    python utils/addComDocFromMQ.py --once      # 只消费一条消息后退出
"""

from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import sys
import hashlib
import tempfile
import time
import logging
from datetime import date, datetime
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import urlopen

import pika

# 让脚本无论从哪里运行都能找到 baseAtoml/llmdeepseek
sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "baseAtoml"))
from llmdeepseek import call_llm

# ---- 日志 ----
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
)
logger = logging.getLogger("addComDocFromMQ")

# ---- 配置 ----
REPO_ROOT = Path(__file__).resolve().parent.parent
MAINDIR = REPO_ROOT / "maindir"

# MQ 配置从 env 读取（延迟加载，避免 import 时 APP_ENV 未设导致崩溃）
sys.path.insert(0, str(REPO_ROOT))
from utils.env_config import load_mq_config, load_mq_queue_name


def _get_mq_url() -> str:
    """获取 RabbitMQ AMQP URL（根据 APP_ENV 自动选 test/uat 配置）。"""
    try:
        return load_mq_config()
    except RuntimeError as e:
        logger.error(f"MQ 配置读取失败: {e}")
        logger.error("请设置 APP_ENV 环境变量，例如: APP_ENV=test python utils/addComDocFromMQ.py --once")
        raise


def _get_mq_queue() -> str:
    """获取 MQ 队列名。"""
    return load_mq_queue_name()

WIKI_MODEL = "deepseek-chat"
WIKI_MAX_TOKENS = 131072  # wiki 生成内容较多，给足够输出空间

# ---- 工具函数 ----


def sha256(text: str, truncate: int = 0) -> str:
    h = hashlib.sha256(text.encode()).hexdigest()
    return h[:truncate] if truncate else h


def parse_json_from_response(text: str) -> dict:
    """从 LLM 响应中提取最外层 JSON 对象。"""
    text = re.sub(r"^```(?:json)?\s*", "", text.strip())
    text = re.sub(r"\s*```$", "", text.strip())
    match = re.search(r"\{[\s\S]*\}", text)
    if not match:
        raise ValueError("响应中未找到 JSON 对象")
    return json.loads(match.group())


# ---- 文档下载与格式转换 ----

RAW_DIR = MAINDIR / "raw"


def download_doc_to_raw(companyId: str, docAddress: str) -> str:
    """
    下载 docAddress 指向的文件到 maindir/raw/<companyId>/<时间戳>.<ext>。
    成功返回本地路径，失败返回原 docAddress（不中断流程）。
    """
    if not docAddress:
        return docAddress

    parsed = urlparse(docAddress)
    # 只对 URL 下载，本地路径直接返回
    if parsed.scheme not in ("http", "https"):
        return docAddress

    path_part = parsed.path or ""
    ext = Path(path_part).suffix.lower()
    if not ext:
        ext = ".bin"

    ts = datetime.now().strftime("%Y%m%d%H%M%S")
    filename = f"{ts}{ext}"

    target_dir = RAW_DIR / companyId
    target_dir.mkdir(parents=True, exist_ok=True)
    target_path = target_dir / filename

    try:
        logger.info(f"  下载文档: {docAddress}")
        req = urlopen(docAddress, timeout=60)
        target_path.write_bytes(req.read())
        logger.info(f"  已保存: {target_path} ({target_path.stat().st_size} bytes)")
        return str(target_path)
    except Exception as e:
        logger.warning(f"  下载失败，使用原始 URL: {e}")
        return docAddress


# 可直接当文本读取的扩展名
TEXT_EXTENSIONS = {".md", ".txt", ".csv", ".json", ".xml", ".yaml", ".yml", ".rst", ".tsv"}

# 需要用 markitdown / pdftotext 转换的扩展名
CONVERTIBLE_EXTENSIONS = {
    ".pdf", ".doc", ".docx", ".pptx", ".xlsx", ".xls",
    ".html", ".htm", ".epub", ".rtf", ".ipynb",
}


def _is_text_usable(text: str) -> bool:
    """检测文本是否有效（排除二进制垃圾）。

    返回 False 如果文本看起来像是原始二进制内容。
    """
    if not text or len(text.strip()) < 10:
        return False

    # null 字节是二进制文件的明显标志
    if "\x00" in text:
        return False

    # 检测常见二进制文件头 magic bytes（取前 200 字符采样）
    head = text[:200]
    _binary_markers = [
        "%PDF-",           # PDF
        "PK\x03\x04",      # DOCX/XLSX/PPTX (ZIP)
        "\xd0\xcf\x11\xe0",  # DOC/XLS (OLE2)
        "\x89PNG",         # PNG
        "\xff\xd8\xff",    # JPEG
    ]
    for marker in _binary_markers:
        if marker in head:
            return False

    # 统计控制字符比例（ASCII 0x00-0x08, 0x0B-0x0C, 0x0E-0x1F, 0x7F）
    # 注意：\n(0x0A), \r(0x0D), \t(0x09) 是正常空白，排除在外
    control_chars = sum(
        1 for c in text
        if ord(c) < 32 and c not in "\n\r\t"
    )
    # 也检测 0x7F-0x9F (DEL + C1 control chars)
    high_control = sum(1 for c in text if 0x7F <= ord(c) <= 0x9F)

    total = len(text)
    bad_ratio = (control_chars + high_control) / total

    # 超过 15% 控制字符 → 判定为二进制垃圾
    if bad_ratio > 0.15:
        return False

    return True


def _convert_file_to_text(file_path: Path) -> str:
    """将任意文件转为文本（markdown 格式）。

    转换策略（按优先级降级）：
      PDF  → pdftotext → PyPDF2 → markitdown
      DOCX → python-docx → markitdown
      XLSX → openpyxl → markitdown
      XLS  → xlrd → markitdown
      DOC  → markitdown（无纯 Python 备选）
      PPTX → python-pptx → markitdown
      其他 → markitdown → 直接读文本

    每层转换后都会用 _is_text_usable 验证结果，无效则降级。
    """
    ext = file_path.suffix.lower()

    # ---- 已是文本格式，直接读 ----
    if ext in TEXT_EXTENSIONS:
        for encoding in ("utf-8", "gbk", "latin-1"):
            try:
                text = file_path.read_text(encoding=encoding)
                if _is_text_usable(text):
                    return text
            except (UnicodeDecodeError, Exception):
                continue
        logger.warning(f"  文本读取失败: {file_path.name}")
        return ""

    # ---- PDF ----
    if ext == ".pdf":
        # 策略 1: pdftotext（系统工具，速度最快，质量最好）
        text = _pdf_via_pdftotext(file_path)
        if text:
            return text

        # 策略 2: PyPDF2（纯 Python，无需系统依赖）
        text = _pdf_via_pypdf2(file_path)
        if text:
            return text

        # 策略 3: markitdown 兜底
        text = _via_markitdown(file_path)
        if text:
            return text

        logger.warning(f"  所有 PDF 解析方式均失败: {file_path.name}")
        return ""

    # ---- DOCX ----
    if ext == ".docx":
        # 策略 1: python-docx
        text = _docx_via_python_docx(file_path)
        if text:
            return text

        text = _via_markitdown(file_path)
        if text:
            return text

        logger.warning(f"  所有 DOCX 解析方式均失败: {file_path.name}")
        return ""

    # ---- XLSX ----
    if ext == ".xlsx":
        text = _xlsx_via_openpyxl(file_path)
        if text:
            return text

        text = _via_markitdown(file_path)
        if text:
            return text

        logger.warning(f"  所有 XLSX 解析方式均失败: {file_path.name}")
        return ""

    # ---- XLS (旧格式) ----
    if ext == ".xls":
        text = _xls_via_xlrd(file_path)
        if text:
            return text

        text = _via_markitdown(file_path)
        if text:
            return text

        logger.warning(f"  所有 XLS 解析方式均失败: {file_path.name}")
        return ""

    # ---- PPTX ----
    if ext == ".pptx":
        text = _pptx_via_python_pptx(file_path)
        if text:
            return text

        text = _via_markitdown(file_path)
        if text:
            return text

        logger.warning(f"  所有 PPTX 解析方式均失败: {file_path.name}")
        return ""

    # ---- DOC / HTML / EPUB / RTF 等 → markitdown 兜底 ----
    text = _via_markitdown(file_path)
    if text:
        return text

    # ---- 最后尝试直接读文本（验证非二进制垃圾） ----
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            text = file_path.read_text(encoding=encoding)
            if _is_text_usable(text):
                logger.info(f"  直接文本读取成功: {file_path.name} ({encoding})")
                return text
        except (UnicodeDecodeError, Exception):
            continue

    logger.warning(f"  无法读取/转换文件: {file_path.name}")
    return ""


# ---- 各格式的独立解析函数 ----


def _pdf_via_pdftotext(file_path: Path) -> str | None:
    """使用 pdftotext 系统工具解析 PDF。"""
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", str(file_path), "-"],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode == 0 and _is_text_usable(result.stdout):
            logger.info(f"  pdftotext 转换成功: {file_path.name} ({len(result.stdout)} 字符)")
            return result.stdout
    except FileNotFoundError:
        pass
    except Exception as e:
        logger.info(f"  pdftotext 失败: {e}")
    return None


def _pdf_via_pypdf2(file_path: Path) -> str | None:
    """使用 PyPDF2 解析 PDF（纯 Python，无需系统依赖）。"""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(file_path))
        parts = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                parts.append(page_text)
        if parts:
            text = "\n\n".join(parts)
            if _is_text_usable(text):
                logger.info(f"  PyPDF2 转换成功: {file_path.name} ({len(text)} 字符)")
                return text
    except ImportError:
        logger.info(f"  PyPDF2 未安装")
    except Exception as e:
        logger.info(f"  PyPDF2 失败: {e}")
    return None


def _docx_via_python_docx(file_path: Path) -> str | None:
    """使用 python-docx 解析 DOCX。"""
    try:
        import docx
        doc = docx.Document(str(file_path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                if row_text.strip():
                    parts.append(row_text)
        if parts:
            text = "\n\n".join(parts)
            if _is_text_usable(text):
                logger.info(f"  python-docx 转换成功: {file_path.name} ({len(text)} 字符)")
                return text
    except ImportError:
        logger.info(f"  python-docx 未安装")
    except Exception as e:
        logger.info(f"  python-docx 失败: {e}")
    return None


def _xlsx_via_openpyxl(file_path: Path) -> str | None:
    """使用 openpyxl 解析 XLSX。"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    parts.append(row_text)
        wb.close()
        if parts:
            text = "\n".join(parts)
            if _is_text_usable(text):
                logger.info(f"  openpyxl 转换成功: {file_path.name} ({len(text)} 字符)")
                return text
    except ImportError:
        logger.info(f"  openpyxl 未安装")
    except Exception as e:
        logger.info(f"  openpyxl 失败: {e}")
    return None


def _xls_via_xlrd(file_path: Path) -> str | None:
    """使用 xlrd 解析旧格式 XLS。"""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(file_path))
        parts = []
        for sheet in wb.sheets():
            parts.append(f"## Sheet: {sheet.name}")
            for row_idx in range(sheet.nrows):
                row_values = sheet.row_values(row_idx)
                row_text = " | ".join(str(cell) for cell in row_values if cell != "")
                if row_text.strip():
                    parts.append(row_text)
        if parts:
            text = "\n".join(parts)
            if _is_text_usable(text):
                logger.info(f"  xlrd 转换成功: {file_path.name} ({len(text)} 字符)")
                return text
    except ImportError:
        logger.info(f"  xlrd 未安装")
    except Exception as e:
        logger.info(f"  xlrd 失败: {e}")
    return None


def _pptx_via_python_pptx(file_path: Path) -> str | None:
    """使用 python-pptx 解析 PPTX。"""
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_parts.append(para.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                        if row_text.strip():
                            slide_parts.append(row_text)
            if slide_parts:
                parts.append(f"## Slide {slide_num}\n" + "\n".join(slide_parts))
        if parts:
            text = "\n\n".join(parts)
            if _is_text_usable(text):
                logger.info(f"  python-pptx 转换成功: {file_path.name} ({len(text)} 字符)")
                return text
    except ImportError:
        logger.info(f"  python-pptx 未安装")
    except Exception as e:
        logger.info(f"  python-pptx 失败: {e}")
    return None


def _via_markitdown(file_path: Path) -> str | None:
    """使用 markitdown 通用转换。"""
    try:
        from markitdown import MarkItDown
        md = MarkItDown(enable_plugins=False)
        result = md.convert(str(file_path))
        if _is_text_usable(result.text_content):
            logger.info(f"  markitdown 转换成功: {file_path.name} ({len(result.text_content)} 字符)")
            return result.text_content
    except ImportError:
        logger.info(f"  markitdown 未安装")
    except Exception as e:
        logger.info(f"  markitdown 失败: {e}")
    return None


def _fetch_from_url(url: str, company_dir: Path) -> str | None:
    """从 HTTP/HTTPS URL 下载文档并转为文本。"""
    try:
        import requests
    except ImportError:
        logger.error("未安装 requests，无法下载文档。pip install requests")
        return None

    parsed = urlparse(url)
    filename = Path(parsed.path).name or "downloaded_doc"
    download_dir = company_dir / "_downloads"
    download_dir.mkdir(parents=True, exist_ok=True)
    local_path = download_dir / filename

    try:
        logger.info(f"  下载文档: {url}")
        resp = requests.get(url, timeout=60, headers={"User-Agent": "addComDocFromMQ/1.0"})
        resp.raise_for_status()

        local_path.write_bytes(resp.content)
        logger.info(f"  已下载: {local_path.name} ({len(resp.content)} bytes)")
    except Exception as e:
        logger.error(f"  下载失败: {url}: {e}")
        return None

    return _convert_file_to_text(local_path)


def _fetch_from_local(doc_address: str) -> str | None:
    """从本地文件或目录读取文档内容。

    - 目录：递归遍历所有可转换文件，拼接内容。
    - 单文件：直接读取/转换。
    """
    local_path = Path(doc_address)
    if not local_path.is_absolute():
        local_path = REPO_ROOT / doc_address

    if not local_path.exists():
        logger.warning(f"  本地路径不存在: {local_path}")
        return None

    # --- 目录：遍历所有文件 ---
    if local_path.is_dir():
        logger.info(f"  检测到目录，遍历所有文件: {local_path}")
        parts = []
        for f in sorted(local_path.rglob("*")):
            if not f.is_file():
                continue
            ext = f.suffix.lower()
            if ext not in TEXT_EXTENSIONS and ext not in CONVERTIBLE_EXTENSIONS:
                logger.info(f"    跳过不支持格式: {f.name}")
                continue
            text = _convert_file_to_text(f)
            if text.strip():
                parts.append(f"\n\n=== 文件: {f.relative_to(local_path)} ===\n\n{text}")
        if not parts:
            logger.warning(f"  目录下无可用文档: {local_path}")
            return None
        return "\n".join(parts)

    # --- 单文件 ---
    return _convert_file_to_text(local_path)


def _fetch_from_oss(oss_url: str, company_dir: Path) -> str | None:
    """从 OSS URL (oss://bucket/path/to/file) 下载文档。需要 oss2 库。"""
    try:
        import oss2
    except ImportError:
        logger.error("未安装 oss2，无法读取 OSS 文档。pip install oss2")
        return None

    # 解析 oss://bucket/path/to/file
    parsed = urlparse(oss_url)
    bucket_name = parsed.hostname or ""
    object_key = parsed.path.lstrip("/")

    if not bucket_name or not object_key:
        logger.error(f"  OSS URL 格式错误: {oss_url}（应为 oss://bucket/path）")
        return None

    # 读取 OSS 配置（兼容 env/oss 或环境变量）
    oss_config = _read_oss_config()
    if not oss_config:
        logger.error("  未找到 OSS 配置（env/oss 或环境变量 OSS_*）")
        return None

    endpoint = oss_config.get("endpoint", "")
    access_key_id = oss_config.get("access_key_id", "")
    access_key_secret = oss_config.get("access_key_secret", "")

    if not all([endpoint, access_key_id, access_key_secret]):
        logger.error("  OSS 配置不完整，需要 endpoint / access_key_id / access_key_secret")
        return None

    filename = Path(object_key).name or "oss_doc"
    download_dir = company_dir / "_downloads"
    download_dir.mkdir(parents=True, exist_ok=True)
    local_path = download_dir / filename

    try:
        logger.info(f"  从 OSS 下载: bucket={bucket_name}, key={object_key}")
        auth = oss2.Auth(access_key_id, access_key_secret)
        bucket = oss2.Bucket(auth, endpoint, bucket_name)
        bucket.get_object_to_file(object_key, str(local_path))
        logger.info(f"  已下载: {local_path.name} ({local_path.stat().st_size} bytes)")
    except Exception as e:
        logger.error(f"  OSS 下载失败: {e}")
        return None

    return _convert_file_to_text(local_path)


def _read_oss_config() -> dict[str, str] | None:
    """读取 OSS 配置：优先 env/oss 文件，其次环境变量。"""
    # 1. env/oss 文件
    oss_file = REPO_ROOT / "env" / "oss"
    if oss_file.exists():
        cfg: dict[str, str] = {}
        for line in oss_file.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                k, v = line.split("=", 1)
                cfg[k.strip()] = v.strip()
        if cfg:
            return cfg

    # 2. 环境变量
    env_keys = {
        "OSS_ENDPOINT": "endpoint",
        "OSS_ACCESS_KEY_ID": "access_key_id",
        "OSS_ACCESS_KEY_SECRET": "access_key_secret",
    }
    cfg = {}
    for env_key, cfg_key in env_keys.items():
        val = os.environ.get(env_key, "")
        if val:
            cfg[cfg_key] = val
    return cfg or None


def fetch_document(doc_address: str, company_dir: Path | None = None) -> str | None:
    """
    从 docAddress 拉取文档内容并转为文本。

    支持来源：
      - HTTP/HTTPS URL → 下载后自动转换（PDF/DOCX/PPTX 等）
      - OSS URL (oss://bucket/path) → 通过 oss2 SDK 下载后转换
      - 本地单文件 → 读取/转换
      - 本地目录 → 递归遍历所有可转换文件，拼接内容

    Args:
        doc_address:  文档地址
        company_dir:  企业 wiki 目录（用于缓存下载文件）

    Returns:
        文档文本，失败返回 None。
    """
    if not doc_address:
        return None

    if company_dir is None:
        company_dir = MAINDIR / "_tmp"

    parsed = urlparse(doc_address)

    if parsed.scheme == "oss":
        return _fetch_from_oss(doc_address, company_dir)

    if parsed.scheme in ("http", "https"):
        return _fetch_from_url(doc_address, company_dir)

    return _fetch_from_local(doc_address)


_FILE_SEP_RE = re.compile(r"\n\n=== 文件: (.+?) ===\n\n")


def _parse_file_list(dc):
    if not dc: return []
    p = _FILE_SEP_RE.split(dc)
    fs = []
    s = 1 if not p[0].strip() else 0
    i = s
    while i + 1 < len(p):
        fn, ft = p[i].strip(), p[i + 1].strip()
        if fn and ft: fs.append((fn, ft))
        i += 2
    if not fs and dc.strip(): fs = [("uploaded", dc.strip())]
    return fs

def build_wiki_prompt(company_id: str, company_name: str, doc_content: str | None, today: str) -> str:
    """构建 wiki 生成的 LLM prompt。

    多文件时：解析文件清单 + 公平分配 512K 字符预算 + 强制要求覆盖所有文件。
    """
    DOC_MAX_CHARS = 512000

    doc_section = ""
    file_count = 1
    if doc_content:
        files = _parse_file_list(doc_content)
        file_count = len(files)
        if file_count <= 1:
            doc_trimmed = doc_content[:DOC_MAX_CHARS]
            doc_section = f"""
企业文档内容（单个文件，最多截取 {DOC_MAX_CHARS} 字符）：
=== DOC START ===
{doc_trimmed}
=== DOC END ===
"""
        else:
            per_file_budget = max(3000, DOC_MAX_CHARS // file_count)

            manifest_lines = []
            for idx, (fname, ftext) in enumerate(files, 1):
                preview = ftext[:100].replace("\n", " ")[:80]
                manifest_lines.append(
                    f"  {idx}. **{fname}** — {len(ftext)} 字符 — 预览: {preview}..."
                )
            manifest = "\n".join(manifest_lines)

            content_parts = []
            for fname, ftext in files:
                trimmed = ftext[:per_file_budget]
                content_parts.append(
                    f"\n\n=== 文件: {fname} ===\n{trimmed}"
                )

            doc_section = f"""
本次共上传 {file_count} 个文件。overview.md 和 index.md 必须覆盖【全部】文件，不得遗漏。

===== 文件清单（{file_count} 个）=====
{manifest}
===== 文件清单结束 =====

以下为各文件内容（每个文件最多截取 {per_file_budget} 字符）：
{"".join(content_parts)}
"""
    else:
        doc_section = """
（未提供文档内容，请根据企业名称和常识生成基础 wiki。）
"""

    file_count_hint = f"{file_count} 个" if doc_content and file_count > 1 else ""

    return f"""你是一个企业 Wiki 知识库构建专家。请为以下企业创建完整的 wiki 知识库。

企业 ID: {company_id}
企业名称: {company_name}
当前日期: {today}
{doc_section}

请根据以上信息，生成以下 wiki 结构的内容。返回纯 JSON（不要 markdown 代码围栏）：

{{
  "index": "完整的 index.md 内容（markdown）。必须在「来源文档」节列出所有 source_page 的链接。",
  "overview": "完整的 overview.md 内容（markdown）。必须综合所有上传文件的内容，覆盖每个文件涉及的核心主题。如果文件较多，为每个文件至少写一句话概括其主题。必须使用 [[页名]] 创建 wiki 内部链接。",
  "source_pages": [
    {{
      "slug": "有意义的英文 kebab-case slug，必须从文件名提炼（如 2026-q1-large-cd、product-manual-v2），禁止使用 source1/source2/doc1 等无意义通用名",
      "title": "文档标题（尽量使用原文件名）",
      "content": "完整的 sources/<slug>.md 内容（markdown）。详细记录文档中的知识。关键术语用 [[页名]] 链接。"
    }}
  ],
  "entity_pages": [
    {{
      "slug": "实体的英文或拼音 kebab-case slug（如 dong-xi-miao、personal-cd-product），禁止使用 entity1/entity2 等通用名",
      "title": "实体标题",
      "content": "完整的 entities/<slug>.md 内容（markdown）。人物、产品、部门、分支机构等实体信息。"
    }}
  ],
  "concept_pages": [
    {{
      "slug": "概念的英文 kebab-case slug（如 net-interest-margin、duration-mismatch），禁止使用 concept1/concept2 等通用名",
      "title": "概念标题",
      "content": "完整的 concepts/<slug>.md 内容（markdown）。行业术语、技术概念、商业模式等。"
    }}
  ],
  "log_entry": "## [{today}] 初始化 | {company_name}\n\n从 {file_count_hint}文档创建企业 wiki。关键信息：..."
}}

要求：
- index.md 必须包含指向所有生成页面的链接，链接路径必须与对应页面的 slug 完全一致
- overview.md 必须覆盖【所有上传文件】的内容，不能遗漏任何文件。生成后请自查：文件清单中的每个文件主题是否都在 overview.md 中有体现
- 每个上传的文件都必须生成对应的 source_page（source_pages 数组长度应等于文件数量）
- 所有 slug 必须是有意义的英文/拼音 kebab-case，严禁使用 source1、entity1、concept1 等无意义通用编号
- 所有页面使用 markdown 格式，大量使用 [[wikilink]] 链接相关页面
- entity_pages 和 concept_pages 可以为空列表 []
- 如果没有文档，至少生成 index.md 和 overview.md
- 只返回 JSON，不要任何解释文字
"""


# 匹配 source1/source2/entity1/concept1/doc1 等无意义通用 slug
_GENERIC_SLUG_RE = re.compile(
    r"^(source|doc|entity|concept|page|file|document)\d*$",
    re.IGNORECASE,
)


def _sanitize_slug(slug: str, title: str, fallback_prefix: str) -> str:
    """检测并替换无意义通用 slug。

    如果 slug 是 source1/doc2/entity1 等无意义命名，则尝试从 title
    生成有意义的 kebab-case slug；若 title 为空或无法处理，则用
    fallback_prefix + 内容哈希生成唯一名。
    """
    if not _GENERIC_SLUG_RE.match(slug):
        return slug

    # 尝试从 title 提取英文/拼音 slug
    if title and title.strip():
        # 中文标题：取前几个字做拼音风格 slug（直接用字符，保留可读性）
        clean = title.strip()
        # 如果标题含英文，提取英文部分做 kebab-case
        eng_parts = re.findall(r"[a-zA-Z0-9]+", clean)
        if eng_parts:
            derived = "-".join(p.lower() for p in eng_parts if len(p) > 1)
            if derived and len(derived) >= 3:
                logger.info(f"  slug 兜底重命名: {slug} → {derived}")
                return derived
        # 纯中文标题：用 hash 前8位保证唯一性
        derived = f"{fallback_prefix}-{sha256(clean, 8)}"
        logger.info(f"  slug 兜底重命名: {slug} → {derived}")
        return derived

    # 无 title，用 fallback_prefix + hash
    derived = f"{fallback_prefix}-{sha256(slug, 8)}"
    logger.info(f"  slug 兜底重命名: {slug} → {derived}")
    return derived


def _ensure_completeness(company_dir: Path, company_name: str) -> None:
    """扫描生成的 wiki 文件，用代码强制保证 index.md 和 overview.md 覆盖所有文件。

    不依赖 LLM 枚举 — 直接从磁盘读取实际生成的文件列表，确保零遗漏。
    """
    index_path = company_dir / "index.md"
    overview_path = company_dir / "overview.md"

    # ── 扫描各目录 ──
    def _scan_md(subdir: str) -> list[tuple[str, str]]:
        """返回 [(文件名, 标题), ...]，按文件名排序。"""
        d = company_dir / subdir
        if not d.is_dir():
            return []
        results = []
        for f in sorted(d.glob("*.md")):
            # 尝试读取第一行 # 标题
            try:
                first_line = f.read_text(encoding="utf-8").split("\n")[0].strip()
                title = first_line.lstrip("#").strip() if first_line.startswith("#") else f.stem
            except Exception:
                title = f.stem
            rel = f"{subdir}/{f.name}"
            results.append((rel, title))
        return results

    sources = _scan_md("sources")
    entities = _scan_md("entities")
    concepts = _scan_md("concepts")

    # ── 重建 index.md ──
    # 保留 LLM 生成的开头描述（从 # 到第一个 ## 之间的内容），其余用代码重建
    old_index = ""
    if index_path.exists():
        old_index = index_path.read_text(encoding="utf-8")

    # 提取开头描述（第一个 ## 之前的内容）
    desc_lines = []
    for line in old_index.split("\n"):
        if line.startswith("##") and desc_lines:
            break
        desc_lines.append(line)
    header = "\n".join(desc_lines).strip()

    # 重建完整 index.md
    lines = [header, ""]
    if not header:
        lines = [f"# {company_name} Wiki", "", "> 企业知识库，代码自动维护文件索引。", ""]
        lines.append("## 概况")
        lines.append("")
        lines.append("- [企业概述](overview.md)")

    if sources:
        lines.append("")
        lines.append("## 来源文档")
        lines.append("")
        for rel, title in sources:
            lines.append(f"- [{title}]({rel})")

    if entities:
        lines.append("")
        lines.append("## 实体")
        lines.append("")
        for rel, title in entities:
            lines.append(f"- [{title}]({rel})")

    if concepts:
        lines.append("")
        lines.append("## 概念")
        lines.append("")
        for rel, title in concepts:
            lines.append(f"- [{title}]({rel})")

    index_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    logger.info(f"  _ensure_completeness: index.md 已重建 ({len(sources)} sources, {len(entities)} entities, {len(concepts)} concepts)")

    # ── 追加 overview.md 覆盖清单 ──
    if sources and overview_path.exists():
        overview = overview_path.read_text(encoding="utf-8")

        # 避免重复追加
        if "## 文档覆盖清单" not in overview:
            coverage = ["", "## 文档覆盖清单", "", "> 以下由代码自动生成，确保所有上传文档均被索引。", ""]
            for rel, title in sources:
                coverage.append(f"- [[{title}|{rel}]]")
            overview_path.write_text(overview.rstrip() + "\n" + "\n".join(coverage) + "\n", encoding="utf-8")
            logger.info(f"  _ensure_completeness: overview.md 已追加覆盖清单 ({len(sources)} 文件)")

    logger.info(f"  _ensure_completeness: 完成")

def create_company_wiki(company_id: str, company_name: str, doc_address: str) -> bool:
    """
    为企业创建 wiki 目录结构和内容。

    Args:
        company_id: 企业 ID（用作文件夹名）
        company_name: 企业名称
        doc_address: 文档地址（URL 或本地路径）

    Returns:
        True 成功，False 失败
    """
    company_dir = MAINDIR / company_id
    company_dir.mkdir(parents=True, exist_ok=True)

    logger.info(f"开始为企业 '{company_name}' ({company_id}) 创建 wiki")
    logger.info(f"  目录: {company_dir}")

    # 0. 下载文档到 maindir/raw/<companyId>/<时间戳>.<ext>
    doc_address = download_doc_to_raw(company_id, doc_address)
    logger.info(f"  文档路径: {doc_address}")

    # 1. 拉取文档内容
    doc_content = fetch_document(doc_address, company_dir)
    if doc_content:
        logger.info(f"  文档已拉取，长度 {len(doc_content)} 字符")
    else:
        logger.warning(f"  未获取到文档内容，将生成基础 wiki")

    # 2. 调用 LLM 生成 wiki 内容
    today = date.today().isoformat()

    prompt = build_wiki_prompt(company_id, company_name, doc_content, today)
    print('vibecodingwhatfuckdoc_content',doc_content)
    print('vibecodingwhatfuckprompt', prompt)

    logger.info(f"  调用 {WIKI_MODEL} 生成 wiki 内容...")
    try:
        raw = call_llm(prompt, model=WIKI_MODEL, max_tokens=WIKI_MAX_TOKENS)
        logger.info(f"  LLM 返回长度 {len(raw)} 字符")
    except Exception as e:
        logger.error(f"  LLM 调用失败: {e}")
        return False
    print('vibecodingwhatfuckdoc_content',doc_content)
    print('vibecodingwhatfuckprompt', prompt)
    # 3. 解析 JSON
    try:
        data = parse_json_from_response(raw)
    except (ValueError, json.JSONDecodeError) as e:
        logger.error(f"  JSON 解析失败: {e}")
        debug_path = company_dir / "_llm_debug.txt"
        debug_path.write_text(raw, encoding="utf-8")
        logger.error(f"  原始响应已保存到 {debug_path}")
        return False
    print('vibecodingwhatfuck raw',raw)
    print('vibecodingwhatfuck data', data)
    # 4. 写入 wiki 文件
    # index.md
    index_content = data.get("index", "")
    if index_content:
        (company_dir / "index.md").write_text(index_content, encoding="utf-8")
        logger.info(f"  wrote: index.md")

    # overview.md
    overview_content = data.get("overview", "")
    if overview_content:
        (company_dir / "overview.md").write_text(overview_content, encoding="utf-8")
        logger.info(f"  wrote: overview.md")

    # sources/
    sources_dir = company_dir / "sources"
    for sp in data.get("source_pages", []):
        slug = _sanitize_slug(sp.get("slug", "source"), sp.get("title", ""), "source")
        content = sp.get("content", "")
        if content:
            sources_dir.mkdir(parents=True, exist_ok=True)
            (sources_dir / f"{slug}.md").write_text(content, encoding="utf-8")
            logger.info(f"  wrote: sources/{slug}.md")

    # entities/
    entities_dir = company_dir / "entities"
    for ep in data.get("entity_pages", []):
        slug = _sanitize_slug(ep.get("slug", "entity"), ep.get("title", ""), "entity")
        content = ep.get("content", "")
        if content:
            entities_dir.mkdir(parents=True, exist_ok=True)
            (entities_dir / f"{slug}.md").write_text(content, encoding="utf-8")
            logger.info(f"  wrote: entities/{slug}.md")

    # concepts/
    concepts_dir = company_dir / "concepts"
    for cp in data.get("concept_pages", []):
        slug = _sanitize_slug(cp.get("slug", "concept"), cp.get("title", ""), "concept")
        content = cp.get("content", "")
        if content:
            concepts_dir.mkdir(parents=True, exist_ok=True)
            (concepts_dir / f"{slug}.md").write_text(content, encoding="utf-8")
            logger.info(f"  wrote: concepts/{slug}.md")

    # log.md
    log_entry = data.get("log_entry", "")
    log_path = company_dir / "log.md"
    if log_entry:
        log_header = (
            f"# {company_name} Wiki Log\n\n"
            "> 记录 wiki 的重要增补、修订和澄清。\n"
        )
        if log_path.exists():
            existing = log_path.read_text(encoding="utf-8").rstrip()
            log_path.write_text(existing + "\n\n" + log_entry + "\n", encoding="utf-8")
        else:
            log_path.write_text(log_header + "\n" + log_entry + "\n", encoding="utf-8")
        logger.info(f"  wrote: log.md")

    # 4.5 代码级完整性保证：重建 index.md + 追加 overview.md 覆盖清单
    _ensure_completeness(company_dir, company_name)

    # 5. 写入任务元信息
    meta = {
        "companyId": company_id,
        "companyName": company_name,
        "docAddress": doc_address,
        "created_at": datetime.now().isoformat(),
        "doc_content_length": len(doc_content) if doc_content else 0,
    }
    (company_dir / ".wiki_meta.json").write_text(
        json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    logger.info(f"  ✅ wiki 创建完成: {company_dir}")
    return True


def consume_one() -> bool:
    """
    从 RabbitMQ 消费一条消息并处理。
    返回 True 表示消费并处理了消息，False 表示队列为空或处理失败。
    """
    try:
        mq_url = _get_mq_url()
        mq_queue = _get_mq_queue()
        connection = pika.BlockingConnection(pika.URLParameters(mq_url))
        channel = connection.channel()
        channel.queue_declare(queue=mq_queue, durable=True)

        # 确认队列中的消息数量
        queue_state = channel.queue_declare(queue=mq_queue, passive=True)
        msg_count = queue_state.method.message_count
        if msg_count == 0:
            logger.info("队列为空，无消息待处理")
            connection.close()
            return False

        # 拉取一条消息
        method_frame, _header_frame, body = channel.basic_get(
            queue=mq_queue, auto_ack=True
        )

        if not body:
            logger.info("未拉取到消息")
            connection.close()
            return False

        connection.close()

        # 解析消息
        msg = json.loads(body)
        task_id = msg.get("task_id", "unknown")
        user = msg.get("user", "")
        company_id = msg.get("companyId", "")
        company_name = msg.get("companyName", "")
        doc_address = msg.get("docAddress", "")

        logger.info(f"[{task_id}] 消费到消息: companyId={company_id}, companyName={company_name}")

        if not company_id:
            logger.error(f"[{task_id}] companyId 为空，跳过")
            return False

        # 处理：创建 wiki
        success = create_company_wiki(company_id, company_name, doc_address)

        if success:
            logger.info(f"[{task_id}] 处理完成 ✅")
        else:
            logger.error(f"[{task_id}] 处理失败 ❌")

        return True

    except pika.exceptions.AMQPConnectionError as e:
        logger.error(f"RabbitMQ 连接失败: {e}")
        return False
    except Exception as e:
        logger.error(f"消费处理异常: {type(e).__name__}: {e}")
        return False


def run_loop(poll_interval: int = 3):
    """
    持续消费循环。

    Args:
        poll_interval: 队列为空时的等待间隔（秒）
    """
    logger.info(f"addComDocFromMQ 消费者启动")
    logger.info(f"  MQ: {_get_mq_url()}")
    logger.info(f"  队列: {_get_mq_queue()}")
    logger.info(f"  maindir: {MAINDIR}")

    while True:
        try:
            consumed = consume_one()
            if not consumed:
                # 队列为空，等待后重试
                time.sleep(poll_interval)
        except KeyboardInterrupt:
            logger.info("收到中断信号，正在退出...")
            break
        except Exception as e:
            logger.error(f"循环异常: {type(e).__name__}: {e}")
            time.sleep(poll_interval)


# ---- 入口 ----

if __name__ == "__main__":
    if "--once" in sys.argv:
        # 只消费一条
        ok = consume_one()
        sys.exit(0 if ok else 1)
    else:
        run_loop()
